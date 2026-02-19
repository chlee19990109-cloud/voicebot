import streamlit as st
import os
import tempfile
import base64
import re
import json
import time

# [AI & LangChain Libraries]
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
# from langchain.docstore.document import Document
# from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate
from pptx import Presentation
from openai import OpenAI
import graphviz # 시각화 필수 라이브러리

# [Multimedia Libraries]
try:
    from moviepy import VideoFileClip
except ImportError:
    try:
        from moviepy.editor import VideoFileClip
    except ImportError:
        VideoFileClip = None

# ==========================================
# [설정] 페이지 기본 설정
# ==========================================
st.set_page_config(page_title="CampusMind", layout="wide", page_icon="🧠")

# [CSS] 스타일링
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@400;500;700&display=swap');
    
    h1, h2, h3, h4, h5, h6, p, li, label, textarea, input, div { 
        font-family: 'Noto Sans KR', sans-serif !important; 
    }
    .material-icons, .material-symbols-rounded {
        font-family: 'Material Icons' !important;
    }
    .stButton > button {
        font-family: 'Noto Sans KR', sans-serif !important;
        width: 100%;
    }
    .stMarkdown {
        font-family: 'Noto Sans KR', sans-serif !important;
    }
    .stTabs button {
        font-family: 'Noto Sans KR', sans-serif !important;
    }
    .stExpander {
        border: 1px solid #e0e0e0;
        border-radius: 8px;
        background-color: #ffffff;
        margin-bottom: 10px;
    }
    .stTabs [data-baseweb="tab-list"] button[aria-selected="true"] {
        background-color: #e3f2fd;
        border-top: 3px solid #1976d2;
        color: #0d47a1;
        font-weight: bold;
    }
    /* 시각화 차트 중앙 정렬 및 크기 최적화 */
    [data-testid="stGraphvizChart"] svg {
        max-width: 100% !important;
        height: auto !important;
        display: block;
        margin: 0 auto;
    }
</style>
""", unsafe_allow_html=True)

# ==========================================
# [언어 팩 (UI)]
# ==========================================
UI = {
    "Korean": {
        "title": "🧠 CampusMind: 지능형 학습 보조 시스템",
        "credit": "By 이충환",
        "caption": "Architecture: RAG-based LLM Workflow",
        "sidebar_title": "⚙️ 데이터 소스",
        "file_label_lec": "📚 강의 자료 (PDF, PPT, Word, 이미지, 음성, 영상 등)",
        "file_label_prob": "📝 연습 문제 (PDF, Word 등)",
        "apikey": "OpenAI API 키",
        "btn_start": "🚀 분석 시작",
        "tabs": ["📝 핵심 정리", "🎨 시각화", "🃏 플래시카드", "🧩 퀴즈", "🎧 오디오 브리핑", "💬 AI 도우미"],
        "input_topic": "주제 필터 (전부, 전체, 빈칸 시 전체 범위)",
        "ph_topic": "예: '신경망' (전부, 전체, 빈칸 시 전체 범위)",
        "msg_proc": "📥 데이터 처리 중...",
        "msg_ingest": "읽는 중: ",
        "msg_done": "✅ 분석 완료!",
        "msg_err_file": "파일 처리 오류: ",
        "msg_nodata": "데이터 없음.",
        "btn_gen": "생성하기",
        "viz_types": ["Mindmap", "Spider Diagram"],
        "quiz_check": "정답 확인",
        "quiz_correct": "정답입니다! ⭕",
        "quiz_wrong": "오답입니다. ❌",
        "quiz_exp": "해설 보기",
        "target_lang": "Korean",
        "lbl_card_front": "질문",
        "lbl_card_back": "정답",
        "audio_btn": "🎙️ 오디오 브리핑 생성",
        "audio_warn": "먼저 요약을 생성해주세요.",
        "spin_gen": "생성 중...",
        "spin_viz": "구조화 중...",
        "spin_audio": "오디오 합성 중...",
        "err_viz": "렌더링 오류. Graphviz가 설치되어 있는지 확인하세요.",
        "err_viz_debug": "DOT 코드 확인 (디버깅)",
        "chat_ph": "질문을 입력하세요...",
        "h_bullet": "1. 핵심 내용 요약",
        "h_table": "2. 상세 요약 표",
        "h_term": "3. 용어 정리",
        "h_th": ["구분", "상세 설명", "용어", "정의", "문맥"],
        "err_json": "데이터 생성 오류. 다시 시도해주세요."
    },
    "English": {
        "title": "🧠 CampusMind: Intelligent Tutor System",
        "credit": "By Choonghwan Lee",
        "caption": "Architecture: RAG-based LLM Workflow",
        "sidebar_title": "⚙️ Data Sources",
        "file_label_lec": "📚 Lecture Materials (PDF, PPT, Word, Image, Audio, Video)",
        "file_label_prob": "📝 Practice Problems (PDF, Word)",
        "apikey": "OpenAI API Key",
        "btn_start": "🚀 Analyze",
        "tabs": ["📝 Summary", "🎨 Visuals", "🃏 Flashcards", "🧩 Quiz", "🎧 Audio Brief", "💬 AI Tutor"],
        "input_topic": "Topic Filter (All, Everything, Blank for All Sections)",
        "ph_topic": "e.g., 'Neural Networks' (All, Everything, Blank for All Sections)",
        "msg_proc": "📥 Processing Data...",
        "msg_ingest": "Ingesting: ",
        "msg_done": "✅ Ready!",
        "msg_err_file": "File Error: ",
        "msg_nodata": "No data.",
        "btn_gen": "Generate",
        "viz_types": ["Mindmap", "Spider Diagram"],
        "quiz_check": "Check Answer",
        "quiz_correct": "Correct! ⭕",
        "quiz_wrong": "Incorrect. ❌",
        "quiz_exp": "Explanation",
        "target_lang": "English",
        "lbl_card_front": "Question",
        "lbl_card_back": "Answer",
        "audio_btn": "🎙️ Generate Audio",
        "audio_warn": "Generate summary first.",
        "spin_gen": "Generating...",
        "spin_viz": "Generating diagram...",
        "spin_audio": "Synthesizing...",
        "err_viz": "Rendering Error. Please check Graphviz installation.",
        "err_viz_debug": "View DOT Code",
        "chat_ph": "Ask a question...",
        "h_bullet": "1. Key Highlights",
        "h_table": "2. Detailed Summary Table",
        "h_term": "3. Terminology",
        "h_th": ["Category", "Detailed Content", "Term", "Definition", "Context"],
        "err_json": "Generation Error. Please try again."
    }
}

# ==========================================
# [Core Logic] 1. Ingestion & Vector DB
# ==========================================
def extract_text(file, ext, path, key):
    if ext == ".pdf":
        return "".join([p.page_content for p in PyPDFLoader(path).load()])
    elif ext in [".docx", ".doc"]:
        return "".join([p.page_content for p in Docx2txtLoader(path).load()])
    elif ext in [".pptx", ".ppt"]:
        prs = Presentation(path)
        return "\n".join([s.text for sl in prs.slides for s in sl.shapes if hasattr(s, "text")])
    elif ext in [".jpg", ".png", ".jpeg"]:
        client = OpenAI(api_key=key)
        with open(path, "rb") as f: enc = base64.b64encode(f.read()).decode('utf-8')
        res = client.chat.completions.create(model="gpt-4o", messages=[{"role": "user", "content": [{"type": "text", "text": "Extract all text visible in this slide/image precisely."}, {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{enc}"}}]}] )
        return f"[Image Source: {file.name}] " + res.choices[0].message.content
    elif ext in [".mp3", ".wav", ".m4a"]:
        client = OpenAI(api_key=key)
        with open(path, "rb") as f: txt = client.audio.transcriptions.create(model="whisper-1", file=f).text
        return f"[Audio Source: {file.name}] " + txt
    elif ext in [".mp4", ".avi", ".mov"]:
        if VideoFileClip is None: return "Error: MoviePy missing."
        audio_path = path + "_temp.mp3"
        try:
            vid = VideoFileClip(path)
            vid.audio.write_audiofile(audio_path, logger=None)
            client = OpenAI(api_key=key)
            with open(audio_path, "rb") as f: txt = client.audio.transcriptions.create(model="whisper-1", file=f).text
            return f"[Video Source: {file.name}] " + txt
        except Exception as e: return str(e)
        finally:
            if os.path.exists(audio_path): os.remove(audio_path)
    return ""

def build_knowledge_base(lec_files, prob_files, key, ui_text):
    docs = []
    status = st.status(ui_text["msg_proc"], expanded=True)
    
    def process_files(file_list, source_type):
        for f in file_list:
            ext = os.path.splitext(f.name)[1].lower()
            status.write(f"{ui_text['msg_ingest']} [{source_type}] {f.name}")
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(f.getvalue())
                tmp_path = tmp.name
            try:
                content = extract_text(f, ext, tmp_path, key)
                if content: 
                    # [핵심] 텍스트 맨 앞에 소스 타입을 명시하여 AI가 구분하기 쉽게 함
                    tagged_content = f"[{source_type}] \n{content}"
                    docs.append(Document(page_content=tagged_content, metadata={"source": f.name, "type": source_type}))
            except Exception as e: st.error(f"{ui_text['msg_err_file']} {e}")
            finally: os.remove(tmp_path)

    if lec_files: process_files(lec_files, "Lecture Material")
    if prob_files: process_files(prob_files, "Practice Problem")
    
    if not docs:
        status.update(label=ui_text["msg_nodata"], state="error")
        return None

    # 청크 사이즈를 조금 늘려서 문맥 파악 능력 향상
    splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    splits = splitter.split_documents(docs)
    db = FAISS.from_documents(splits, OpenAIEmbeddings(openai_api_key=key))
    status.update(label=ui_text["msg_done"], state="complete", expanded=False)
    return db

def get_rag_chain(db, key, target_lang):
    llm = ChatOpenAI(model="gpt-4o", temperature=0.2, openai_api_key=key)
    
    # [핵심] 연습 문제 태그를 보고 스타일을 분석하도록 지시
    template = f"""
    You are an intelligent AI Teaching Assistant and Exam Strategist.
    
    *** SOURCE IDENTIFICATION ***
    - Text starting with `[Lecture Material]` is conceptually explanatory.
    - Text starting with `[Practice Problem]` contains actual exam/quiz questions.

    *** INSTRUCTIONS ***
    1. **Concept Explainer**: If asked about concepts, prioritize `[Lecture Material]`.
    2. **Exam Strategist**: If asked about "exam style", "preparation", or "types of problems":
       - Look strictly at the content labeled `[Practice Problem]`.
       - Analyze the format (Multiple choice? Essay? Calculation?) and difficulty.
       - Provide a strategy based on those specific patterns.
       - If no `[Practice Problem]` is found, state that you need practice files to analyze the exam style.
    
    *** STRICT RULES ***
    - Answer ONLY using the provided [Context].
    - Output Language: **{target_lang}**.
    
    [Context]:
    {{context}}
    
    [User Question]:
    {{question}}
    """
    
    retriever = db.as_retriever(
        search_type="mmr", 
        search_kwargs={'k': 20, 'fetch_k': 50, 'lambda_mult': 0.6}
    )
    
    return RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        chain_type_kwargs={"prompt": PromptTemplate(template=template, input_variables=["context", "question"])}
    )

# ==========================================
# [Core Logic] 2. Generation Functions
# ==========================================
def get_scope(topic):
    return "the ENTIRE provided material" if not topic or not topic.strip() else f"the topic '{topic}'"

def clean_json(text):
    text = text.strip()
    text = re.sub(r"```(json)?", "", text, flags=re.IGNORECASE).replace("```", "")
    match = re.search(r"(\[.*\])", text, re.DOTALL)
    if match: text = match.group(1)
    text = re.sub(r",\s*\]", "]", text)
    return text

def clean_dot_code(text):
    text = text.strip()
    text = re.sub(r"```(dot)?", "", text).replace("```", "")
    start_idx = text.find("digraph")
    if start_idx == -1: return text 
    open_brace = text.find("{", start_idx)
    if open_brace == -1: return text
    close_brace = text.rfind("}")
    if close_brace == -1: return text
    return text[start_idx : close_brace+1]

# 용어 정리 및 요약
def gen_summary(db, api_key, topic, ui_text):
    # 입력값 분석 (전체 vs 특정 토픽)
    is_all_mode = False
    if not topic or topic.strip().lower() in ["all", "전부", "전체", "everything"]:
        is_all_mode = True
        scope_text = "the ENTIRE provided material (All lectures)"
    else:
        scope_text = f"the specific topic '{topic}'"

    lang = ui_text["target_lang"]
    
    # 모드에 따른 설정
    if is_all_mode:
        # [전체 모드]
        # 검색어: 전체 구조를 파악할 수 있는 포괄적 키워드
        search_query = "Table of contents, Lecture titles, Course outline, Key concepts summary"
        k_val = 80  # 전체를 봐야 하므로 많은 청크(80개)를 가져옴
        
        mode_instruction = f"""
        - **Goal**: Create a **"Master Course Outline"** that lists **EVERY** detected file or lecture.
        - **Constraint**: Keep descriptions concise to ensure ALL lectures are covered within the output limit.
        - **Coverage**: It is critical to list **ALL** lectures/files found in the text. Do not stop after the first few.
        - **Format**: For each lecture, provide a brief summary and a list of key exam concepts.
        """
        
        # [전체 모드 가이드라인]
        guidelines = f"""
        1. **Context-Based**: Answer ONLY based on the provided [Context].
        2. **Completeness (CRITICAL)**: 
           - You MUST iterate through **ALL** detected files/lectures.
           - Do not skip the later lectures. 
        3. **Terminology Integrity (STRICT)**: 
           - Even in the concept list, terms must be a **VERBATIM COPY** from the source.
           - **DO NOT TRANSLATE THE TERM ITSELF.**
        """

        # [전체 모드 포맷] (표 없는 경량화 구조)
        format_instruction = f"""
        **[INSTRUCTION: Repeat the block below for EVERY detected File/Lecture]**

        ## 📂 [Insert File or Lecture Name]
        
        ### 📖 {ui_text['h_bullet']} (Lecture Overview)
        - (Summarize the main theme of this lecture in 2-3 sentences in {lang}.)
        
        ### 🔑 Key Exam Concepts
        - **(Concept 1)**: (Short definition/Core logic in {lang})
        - **(Concept 2)**: (Short definition/Core logic in {lang})
        - **(Concept 3)**: (Short definition/Core logic in {lang})
        
        ---
        """
    
    else:
        # [특정 토픽 모드]
        search_query = topic
        k_val = 15  # 특정 토픽 집중
        
        mode_instruction = f"""
        - **Scope Focus**: Focus **DEEPLY and STRICTLY** on the concept of '{topic}'. Ignore unrelated sections.
        - **Terminology Strategy**: Select terms that are **semantically related** to '{topic}' (e.g., sub-concepts, components, algorithms).
        """
        
        # [특정 토픽 가이드라인] (제공해주신 내용 그대로 적용)
        guidelines = f"""
        1. **Context-Based**: Answer ONLY based on the provided [Context].
        2. **Comprehensive Coverage**: 
           - Do NOT limit the number of key points.
           - Extract **ALL** core concepts, definitions, formulas, and arguments.
           - Aim for high detail.
        3. **Terminology Integrity (STRICT)**: 
           - In the Terminology Table, the 'Term' column must be a **VERBATIM COPY** from the source text.
           - **DO NOT TRANSLATE THE TERM ITSELF.**
           - If the source uses English (e.g., "Backpropagation"), keep it "Backpropagation".
           - If the source uses Korean (e.g., "역전파"), keep it "역전파".
           - Only the Definition and Context columns should be in **{lang}**.
        """
        
        # [특정 토픽 포맷]
        format_instruction = f"""
        ### {ui_text['h_bullet']}
        - (List **ALL** exam-relevant key points about '{topic}' in {lang}.)
        
        ### {ui_text['h_table']}
        | {ui_text['h_th'][0]} | {ui_text['h_th'][1]} |
        |---|---|
        | (Category in {lang}) | (Detailed explanation in {lang}) |
        
        ### {ui_text['h_term']}
        | {ui_text['h_th'][2]} | {ui_text['h_th'][3]} | {ui_text['h_th'][4]} |
        |---|---|---|
        | **(EXACT SOURCE TERM)** | (Definition in {lang}) | (Context/Relation in {lang}) |
        """

    # [핵심] DB 검색 수행
    # DB가 비어있는지 체크
    if db is None:
        return "Error: Database is not initialized."
    
    # DB에서 텍스트 추출
    docs = db.similarity_search(search_query, k=k_val)
    context = "\n".join([d.page_content for d in docs])

    # 프롬프트 조합
    prompt = f"""
    You are an expert **Professor** and **Exam Preparation Tutor**.
    Analyze {scope_text} based STRICTLY on the provided context.
    
    *** MODE INSTRUCTION ***
    {mode_instruction}
    
    *** CRITICAL GUIDELINES ***
    {guidelines}

    *** OUTPUT FORMAT ***
    {format_instruction}

    [Context]:
    {context}
    """
    
    # [핵심] LLM 직접 호출
    try:
        llm = ChatOpenAI(model="gpt-4o", temperature=0.3, openai_api_key=api_key)
        response = llm.invoke(prompt)
        return response.content
    except Exception as e:
        return f"Error during generation: {str(e)}"

# 시각화 (마인드맵, 스파이더 다이어그램)
def gen_diagram_optimized(db, api_key, topic, viz_type, ui_text):
    # 1. 입력값 분석 (전체 vs 특정 토픽)
    is_all_mode = False
    if not topic or topic.strip().lower() in ["all", "전부", "전체", "everything"]:
        is_all_mode = True
        
    # 2. 검색 전략 및 프롬프트 지침 설정
    if is_all_mode:
        # [전체 모드]
        search_query = "Table of Contents, Course Syllabus, All Lecture Titles, All Chapter Titles, Lecture 1, Lecture 2, ..., Lecture N"
        search_type = "mmr" 
        k_val = 300     # 전체 범위를 커버하기 위해 유지
        fetch_k = 3000
        
        root_node = "Course Overview"
        
        # ★ 전체 모드 지침
        scope_instruction = """
        - **MODE**: Full Course Syllabus & Key Concepts.
        - **GOAL**: Visualize **EVERY SINGLE** Lecture/Chapter found in the files, and optionally attach 2-3 key concepts to each lecture.
        - **CRITICAL REQUIREMENT (NO OMISSION)**: 
            1. **Exhaustive List**: Look at the [Source File] names and context. You MUST create a node for EVERY lecture present (e.g., Lecture 1 to Lecture N). **DO NOT SKIP ANY LECTURE.**
            2. **Hierarchy**: Root -> Lecture Node (Level 1) -> Keyword Nodes (Level 2).
            3. **NO EDGE LABELS**: Edges must be plain lines. **Put all text INSIDE the Node.**
            4. **Logical Order**: Arrange nodes in the order.
        - **NAMING RULES (CRITICAL)**:
            1. Node Label: "Lec X: [Title]" (e.g., "Lec 2: Metals", "Lec 5: Composites").
               - BAD: "Lecture 2" -> "Metals" (Do not split).
               - GOOD: Root -> "Lec 2: Metals".
            2. Keyword Label: Use the exact term from the context (e.g., "Thermodynamics", "Stress-Strain Curve").
            3. **NO EDGE LABELS**: Edges must be plain lines. Text goes inside nodes.
        """
        
    else:
        # [특정 토픽 모드]
        search_query = f"Structure and details of '{topic}', sub-types, components, key features"
        search_type = "similarity" 
        k_val = 15
        fetch_k = 0
        
        root_node = topic.strip()
        
        scope_instruction = f"""
        - **MODE**: Structured Deep Dive.
        - **GOAL**: Visualize the **Structure** of '{topic}' concisely.
        - **STYLE**:
            1. Root ('{topic}') -> Sub-Components / Types (Level 1).
            2. Sub-Components -> Key Characteristics (Level 2).
            3. **Constraint**: Use short phrases in nodes (Max 5-8 words). Avoid long sentences.
            4. **NO EDGE LABELS**: Edges must be plain lines. Text goes inside nodes.
        """

    # 3. 문서 검색
    try:
        if search_type == "mmr":
            # fetch_k를 충분히 주어 다양성 확보
            docs = db.max_marginal_relevance_search(search_query, k=k_val, fetch_k=fetch_k)
        else:
            docs = db.similarity_search(search_query, k=k_val)
    except Exception as e:
        print(f"Search Error: {e}")
        docs = db.similarity_search(search_query, k=k_val)
    
    # ★ 핵심 파트: Context에 '파일명(Source)'을 직접 명시하여 AI가 누락 없이 전체 강의를 파악하게 함
    context_chunks = []
    for d in docs:
        # metadata에서 파일명을 추출 (경로 제외)
        source = d.metadata.get('source', '')
        if source:
            filename = source.split('/')[-1].split('\\')[-1]
            context_chunks.append(f"--- [Source File: {filename}] ---\n{d.page_content}")
        else:
            context_chunks.append(d.page_content)
            
    context = "\n\n".join(context_chunks)
    
    if not context:
        return 'digraph G { "No Data" [shape=box]; }'
    
    # Context 길이 제한 (토근 초과 방지)
    # 한글/영어 혼용 시 1토큰 ≈ 2~3 char.
    # 50,000자 ≈ 15,000 ~ 20,000 토큰 (안전 구간)
    safe_context = context[:50000]

    llm = ChatOpenAI(model="gpt-4o", temperature=0, openai_api_key=api_key)
    font_attr = 'fontname="Malgun Gothic, AppleGothic, sans-serif"'
    
    # 4. 시각화 스타일 설정
    no_edge_text = 'label="", xlabel="",' 
    
    if "Mind" in viz_type:
        layout_engine = "dot"
        rank_dir = "LR" 
        
        # 전체 모드일 경우 노드 간격(ranksep)을 조금 더 좁혀서 한눈에 들어오게 조정
        sep_settings = 'nodesep=0.25; ranksep=0.8;' if is_all_mode else 'nodesep=0.3; ranksep=1.0;'
        graph_attr = f'rankdir={rank_dir}; splines=ortho; {sep_settings} compound=true;'
        
        # 노드 스타일
        if is_all_mode:
            # 전체 모드: 박스형, 연한 파랑
            node_def = f'node [shape=box, style="filled,rounded", fillcolor="#E3F2FD", penwidth=1.0, fontsize=12, {font_attr}];'
        else:
            # 상세 모드: 노트형, 연한 노랑
            node_def = f'node [shape=note, style="filled,rounded", fillcolor="#FFF9C4", penwidth=1.0, fontsize=12, margin="0.1,0.1", {font_attr}];'
            
        edge_def = f'edge [arrowhead=vee, arrowsize=0.5, color="#546E7A", {no_edge_text} {font_attr}];'
        
        viz_rules = f"""
        2. **Mindmap Rules**:
            - **Root Node**: Label: **"{root_node}"** (Shape: doubleoctagon, Color: #FFCCBC).
            - **NO EDGE TEXT**: Strictly forbidden. Use plain lines only.
            - **Consistency**: Ensure lecture names correspond to the context provided.
        """
    else:
        # Spider Diagram
        layout_engine = "neato"
        graph_attr = 'overlap=false; splines=curved; sep="+25,25"; esep="+10,10"; start=regular;'
        node_def = 'node [shape=plaintext, fontcolor="#37474F", fontsize=11, ' + font_attr + '];'
        edge_def = f'edge [arrowhead=none, color="#B0BEC5", len=2.5, penwidth=1.0, {no_edge_text} {font_attr}];'
        
        viz_rules = f"""
        2. **Spider Diagram Rules**:
            - **Root Node**: Center node **"{root_node}"**.
            - **NO EDGE TEXT**: Strictly forbidden.
        """

    # 5. 프롬프트 조합
    prompt = f"""
    Role: Expert Curriculum Designer & Data Visualization Specialist.
    Task: Generate Graphviz DOT code based on the [Context].
    
    *** VISUALIZATION INSTRUCTION ***
    {scope_instruction}
    
    [Context]
    {safe_context} 
    
    *** STRICT RULES ***
    1. Use ONLY information from the Context.
    2. **Language**: Use the same language as the Context.
    3. **CLEAN EDGES**: **NEVER** put text on edges. Just A -> B.
    4. **NO GENERIC NAMES**: Use the real lecture titles from the [Source File] names or text.
       - FORBIDDEN: "keyword1", "nodeA", "Lecture X".
       - REQUIRED: "Structure", "Thermodynamics", "Lec 2: Metals".
    {viz_rules}
    
    Template:
    digraph G {{
        layout={layout_engine};
        {graph_attr}
        {node_def}
        {edge_def}
        
        // Root Node
        root [label="{root_node}", shape=doubleoctagon, style=filled, fillcolor="#FFCCBC", fontsize=14];
                
        // Define Nodes & Edges
        // ... (Generate nodes for ALL chapters found in context)
    }}
    """
    
    try:
        # invoke 사용 (LangChain 최신 버전 호환)
        res = llm.invoke(prompt).content
        return clean_dot_code(res)
    except Exception as e:
        # 에러 발생 시 사용자에게 힌트를 주는 노드 생성
        return f'digraph G {{ "Error" [label="Error: {str(e)[:40]}...", shape=box, style=filled, fillcolor="#FFCDD2"]; }}'

# 플래시카드
def gen_flashcards(db, api_key, topic, ui_text):
    lang = ui_text["target_lang"]
    
    # 1. 입력값 분석 (전체 vs 특정 토픽)
    is_all_mode = False
    if not topic or topic.strip().lower() in ["all", "전부", "전체", "everything"]:
        is_all_mode = True

    # 2. 검색 및 프롬프트 전략 설정
    if is_all_mode:
        # [전체 모드]
        search_query = "Important definitions, core concepts, exam questions, summary"
        # k를 대폭 늘림 (더 많은 내용을 참조하여 많이 생성하기 위함)
        k_val = 80 
        scope_instruction = """
        - **Quantity**: **DO NOT LIMIT** the number of cards. Generate as many flashcards as possible to cover the entire context exhaustively.
        - **Scope**: Cover the **ENTIRE breadth** of the provided material from start to finish.
        - **Diversity**: Extract key questions from ALL sections (intro, body, conclusion).
        """
    else:
        # [특정 토픽 모드]
        search_query = topic
        k_val = 15 # 특정 토픽
        scope_instruction = f"""
        - **Quantity**: Create a comprehensive set of flashcards (no fixed limit) to fully master '{topic}'.
        - **Scope**: Focus **STRICTLY** on the concept of '{topic}'.
        - **Depth**: Ask about definitions, sub-concepts, differences, and applications related specifically to '{topic}'.
        """

    # 3. 문서 검색 (늘어난 k_val 사용)
    docs = db.similarity_search(search_query, k=k_val)
    context = "\n".join([d.page_content for d in docs])
    
    # 4. 프롬프트 생성
    client = OpenAI(api_key=api_key)
    prompt = f"""
    Role: Exam Prep Tutor.
    Task: Create a comprehensive list of Q&A flashcards based on the [Context].
    
    *** SCOPE INSTRUCTION ***
    {scope_instruction}

    Language: {lang}.
    Format: JSON Array ONLY. Keys: "front" (Question), "back" (Short Answer).
    
    [Context]:
    {context[:15000]} # 컨텍스트 길이 제한을 15000으로 대폭 늘림 (재료가 많아야 많이 만듦)
    
    Output example: [{{"front": "What is X?", "back": "X is Y."}}, {{"front": "...", "back": "..."}}]
    """
    
    try:
        # max_tokens를 늘리거나 기본값(모델 최대치)을 사용하도록 두어 출력이 잘리지 않게 함
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "system", "content": "You are a JSON generator."}, 
                      {"role": "user", "content": prompt}],
            temperature=0.5
        )
        res = response.choices[0].message.content
        return json.loads(clean_json(res))
    except Exception as e:
        return [{"front": "Error", "back": f"Failed to generate: {str(e)}"}]

# 퀴즈
def gen_quiz(db, api_key, topic, ui_text):
    lang = ui_text["target_lang"]
    
    # 1. 입력값 분석 (전체 vs 특정 토픽)
    is_all_mode = False
    if not topic or topic.strip().lower() in ["all", "전부", "전체", "everything", "total"]:
        is_all_mode = True

    # 2. 검색 및 프롬프트 전략 설정
    if is_all_mode:
        # [전체 모드]
        search_query = "Exam questions, practice problems, core concepts, critical knowledge"
        # 문제를 많이 내려면 재료가 많아야 하므로 k를 80으로 대폭 증가
        k_val = 80
        scope_instruction = """
        - **Quantity**: **DO NOT LIMIT** the number of questions. Generate as many unique questions as possible (e.g., 10, 20, or more) to cover the entire context exhaustively.
        - **Scope**: Questions must cover **various lectures/sections** of the provided material, not just one.
        - **Diversity**: Ensure questions range from fundamental definitions to complex applications found across the entire text.
        """
    else:
        # [특정 토픽 모드]
        search_query = topic
        k_val = 15 # 특정 토픽도 깊게 파기 위해 조금 늘림
        scope_instruction = f"""
        - **Quantity**: Create a comprehensive set of questions (no fixed limit) to fully master '{topic}'.
        - **Scope**: Focus **STRICTLY** on the concept of '{topic}'.
        - **Depth**: Create questions that test the definition, usage, nuances, and common misconceptions of '{topic}' specifically.
        """

    # 3. 문서 검색 (늘어난 k_val 사용)
    docs = db.similarity_search(search_query, k=k_val)
    context = "\n".join([d.page_content for d in docs])
    
    # 4. 프롬프트 생성
    client = OpenAI(api_key=api_key)
    prompt = f"""
    Role: Professor.
    Task: Create a comprehensive set of multiple-choice questions based on [Context].
    
    *** SCOPE INSTRUCTION ***
    {scope_instruction}

    Language: {lang}.
    Format: JSON Array ONLY.
    
    Requirements:
    - 4 Options per question.
    - Include clear "explanation" for the correct answer.
    - **Randomize the position of the correct answer** (do not always make 'A' the answer).
    - **IMPORTANT**: The 'answer' field must be the **EXACT String value** from the 'options' list, NOT just 'A', 'B', 'C', or 'D'.
    
    [Context]:
    {context[:15000]} # 컨텍스트 길이를 15000으로 대폭 늘림 (문제를 많이 만들기 위함)
    
    Output example: 
    [{{"question":"What is 1+1?", "options":["3","2","5","4"], "answer":"2", "explanation":"1+1 equals 2."}}]
    """
    
    try: 
        # max_tokens 제한에 걸리지 않도록 주의 (GPT-4o는 출력 토큰 여유가 큼)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "system", "content": "You are a JSON generator."}, 
                      {"role": "user", "content": prompt}],
            temperature=0.5
        )
        res = response.choices[0].message.content
        return json.loads(clean_json(res))
    except Exception as e:
        return [{"question": "Error", "options": ["Error"], "answer": "Error", "explanation": str(e)}]

# ==========================================
# [UI] Main Application
# ==========================================
keys = ["chain", "summary", "diagram", "quiz_data", "flashcards", "messages", "db", "api_key"]
for k in keys:
    if k not in st.session_state: st.session_state[k] = None
if st.session_state.messages is None: st.session_state.messages = []

with st.sidebar:
    lang_opt = st.radio("언어 모드 / Language Mode", ["Korean", "English"], horizontal=True)
    ui = UI[lang_opt]
    
    st.title(ui["sidebar_title"])
    api_key_input = st.text_input(ui["apikey"], type="password")
    
    lec_files = st.file_uploader(ui["file_label_lec"], accept_multiple_files=True, key="lec")
    prob_files = st.file_uploader(ui["file_label_prob"], accept_multiple_files=True, key="prob")
    
    if st.button(ui["btn_start"], type="primary"):
        if api_key_input and (lec_files or prob_files):
            st.session_state.api_key = api_key_input
            # 분리된 파일 리스트를 DB 생성 함수에 전달
            db = build_knowledge_base(lec_files, prob_files, api_key_input, ui)
            if db:
                st.session_state.db = db
                st.session_state.chain = get_rag_chain(db, api_key_input, ui["target_lang"])
                st.session_state.summary = None
                st.session_state.diagram = None
                st.rerun()

st.title(ui["title"])
st.markdown(f"**{ui['credit']}**")
st.caption(ui["caption"])

if st.session_state.chain and st.session_state.db:
    t1, t2, t3, t4, t5, t6 = st.tabs(ui["tabs"])

    with t1: # 요약
        topic_s = st.text_input("Topic_Sum", placeholder=ui["ph_topic"], label_visibility="collapsed")
        if st.button(ui["btn_gen"], key="sum"):
            with st.spinner(ui["spin_gen"]):
                st.session_state.summary = gen_summary(
                    st.session_state.db,       # DB 객체 전달
                    st.session_state.api_key,  # api_key 전달
                    topic_s, 
                    ui
                    )
        if st.session_state.summary: st.markdown(st.session_state.summary)

    with t2: # 시각화
        c1, c2, c3 = st.columns([2, 3, 1])
        with c1:
            v_type = st.selectbox("Style", ui["viz_types"], label_visibility="collapsed")
        with c2:
            topic_v = st.text_input("Viz_Topic", placeholder=ui["ph_topic"], label_visibility="collapsed")
        with c3:
            if st.button(ui["btn_gen"], key="viz", use_container_width=True):
                with st.spinner(ui["spin_viz"]):
                    st.session_state.diagram = gen_diagram_optimized(
                        st.session_state.db, 
                        st.session_state.api_key, 
                        topic_v, 
                        v_type, 
                        ui
                    )
        
        if st.session_state.diagram:
            try:
                st.graphviz_chart(st.session_state.diagram, use_container_width=True)
                with st.expander(ui["err_viz_debug"]):
                    st.code(st.session_state.diagram, language="dot")
            except Exception as e:
                st.error(f"{ui['err_viz']} ({str(e)})")
                st.code(st.session_state.diagram)

    with t3: # 플래시카드
        topic_f = st.text_input("Topic_Flash", placeholder=ui["ph_topic"], label_visibility="collapsed")
        if st.button(ui["btn_gen"], key="flash"):
            with st.spinner(ui["spin_gen"]):
                # 데이터 생성 및 저장
                st.session_state.flashcards = gen_flashcards(st.session_state.db, st.session_state.api_key, topic_f, ui)
        
        # 화면 출력 로직
        if st.session_state.flashcards:
            # 에러 메시지가 담겨 있는지, 정상 데이터인지 확인
            if isinstance(st.session_state.flashcards, list) and len(st.session_state.flashcards) > 0 and "front" in st.session_state.flashcards[0]:
                cols = st.columns(2)
                for i, c in enumerate(st.session_state.flashcards):
                    with cols[i % 2]:
                        st.info(f"**Q{i+1}: {c['front']}**")
                        with st.expander(ui['lbl_card_back']):
                            st.write(c['back'])
            else:
                # 생성 실패 시 에러 메시지 출력
                st.error(ui["err_json"])
                st.write(st.session_state.flashcards)

    with t4: # 퀴즈
        topic_q = st.text_input("Topic_Quiz", placeholder=ui["ph_topic"], label_visibility="collapsed")
        if st.button(ui["btn_gen"], key="quiz"):
            with st.spinner(ui["spin_gen"]):
                # 데이터 생성 및 저장
                st.session_state.quiz_data = gen_quiz(st.session_state.db, st.session_state.api_key, topic_q, ui)
                # 퀴즈는 라디오 버튼 상태 관리를 위해 rerun이 유용할 수 있음
                st.rerun()

        # 화면 출력 로직
        if st.session_state.quiz_data:
            if isinstance(st.session_state.quiz_data, list) and len(st.session_state.quiz_data) > 0 and "question" in st.session_state.quiz_data[0]:
                for i, q in enumerate(st.session_state.quiz_data):
                    st.markdown(f"#### Q{i+1}. {q['question']}")
                    
                    # 라디오 버튼 (선택지)
                    ans = st.radio(
                        "Select:", 
                        q['options'], 
                        key=f"q_{i}", 
                        index=None, 
                        label_visibility="collapsed"
                    )
                    
                    # 정답 확인 버튼
                    if st.button(ui["quiz_check"], key=f"chk_{i}"):
                        if ans == q['answer']: 
                            st.success(ui["quiz_correct"])
                        else: 
                            st.error(ui["quiz_wrong"])
                        
                        # 해설 보기
                        with st.expander(ui["quiz_exp"]): 
                            st.write(q['explanation'])
                    st.divider()
            else:
                # 생성 실패 시 에러 메시지 출력
                st.error(ui["err_json"])
                st.write(st.session_state.quiz_data)
    
    with t5: # 오디오
        if st.button(ui["btn_gen"], key="audio"):
            # 전제 조건 확인
            if st.session_state.summary:
                client = OpenAI(api_key=st.session_state.api_key)
                
                # [기능 1] 마크다운 기호 제거 함수 (듣기 편하게)
                def clean_markdown_for_speech(text):
                    # 1. 헤더 제거 (### 등)
                    text = re.sub(r'#+\s?', '', text)
                    # 2. 볼드체/이탤릭 제거 (** **)
                    text = re.sub(r'\*\*|__', '', text)
                    # 3. 불필요한 공백/줄바꿈 정리
                    text = re.sub(r'\n+', ' ', text)
                    return text.strip()

                # [기능 2] 모든 자료의 '핵심 내용'만 추출하는 함수 (업그레이드)
                def extract_all_core_parts(text, ui_text):
                    start_marker = ui_text['h_bullet'] # "1. 핵심 내용 요약"
                    end_marker = ui_text['h_table']    # "2. 상세 요약 표"
                    
                    # 정규표현식으로 start와 end 사이의 모든 텍스트 추출 (re.DOTALL: 줄바꿈 포함)
                    # 패턴: (시작마커) ...내용... (끝마커)
                    pattern = f"{re.escape(start_marker)}(.*?){re.escape(end_marker)}"
                    matches = re.findall(pattern, text, re.DOTALL)
                    
                    if matches:
                        # 추출된 모든 섹션을 하나로 합침
                        combined_text = " ".join(matches)
                        return clean_markdown_for_speech(combined_text)
                    else:
                        # 매칭 실패 시 (예외 처리)
                        return clean_markdown_for_speech(text[:1000])

                try:
                    with st.spinner(ui["spin_audio"]):
                        # 1. 텍스트 추출 및 정제
                        core_summary = extract_all_core_parts(st.session_state.summary, ui)
                        
                        # 2. 길이 제한 (OpenAI TTS 한도 4096자 고려, 안전하게 4000자)
                        if len(core_summary) > 4000:
                            final_input = core_summary[:4000] + "... (Content truncated due to length limit)"
                            st.caption("⚠️ 텍스트가 너무 길어 앞부분 4000자만 재생됩니다.")
                        else:
                            final_input = core_summary

                        # 3. TTS 생성
                        audio = client.audio.speech.create(
                            model="tts-1",
                            voice="alloy",
                            input=final_input
                        )
                        
                        st.success("Audio generated! (Reading all 'Key Highlights')")
                        
                        # 4. 재생 및 스크립트 확인
                        st.audio(audio.content, format="audio/mp3")
                        
                        with st.expander("📜 읽어준 대본 (Script)"):
                            st.write(final_input)

                except Exception as e: 
                    st.error(f"Error: {str(e)}")
            else: 
                st.warning(ui["audio_warn"])
    
    with t6: # AI 튜터
        chat_box = st.container(height=500)
        for m in st.session_state.messages: chat_box.chat_message(m["role"]).write(m["content"])
        
        if q := st.chat_input(ui["chat_ph"]):
            # 유저 메시지 표시 및 저장
            st.session_state.messages.append({"role":"user", "content":q})
            chat_box.chat_message("user").write(q)
            
            with st.spinner("Analyzing intent & Searching documents..."):
                try:
                    # === [핵심 로직] 질문 의도 분석 및 쿼리 확장 ===
                    q_lower = q.strip().lower()
                    
                    # 1. [전체 요약 모드] ("전부", "all" 등)
                    if q_lower in ["all", "전부", "전체", "everything", "요약해줘"]:
                        # 검색 효과: 문서 전체를 아우르는 키워드로 검색
                        # 지시 효과: 전체를 상세히 정리하라는 명령 추가
                        search_query = (
                            "Provide a comprehensive and very detailed summary of the ENTIRE provided material. "
                            "Cover all lectures, core concepts, structure, and main arguments from start to finish. "
                            "Do not miss any major sections."
                        )
                    
                    # 2. [시험 전략 모드] ("시험", "유형", "전략", "계획" 등)
                    elif any(x in q_lower for x in ["시험", "exam", "test", "유형", "type", "strategy", "plan", "계획", "대비"]):
                        # 검색 효과: [Practice Problem] 태그가 붙은 내용 위주로 검색 유도
                        # 지시 효과: 스타일 분석 및 전략 수립 요청
                        search_query = (
                            f"User Question: '{q}'\n\n"
                            "Task: Act as an Exam Strategist. "
                            "1. Analyze the content labeled `[Practice Problem]` to identify exam styles (MCQ, Essay, etc.) and difficulty. "
                            "2. Summarize the types of questions that appear. "
                            "3. Provide a concrete study plan and preparation strategy based on these patterns."
                        )
                    
                    # 3. [특정 용어/일반 질문 모드]
                    else:
                        # 검색 효과: 해당 용어와 관련된 문맥 검색
                        # 지시 효과: 단답형이 아닌 '상세 설명' 유도
                        search_query = (
                            f"Explain the concept of '{q}' in great detail. "
                            "Include its definition, context, related terms, and why it is important in this document."
                        )
                    
                    # RAG 체인 실행 (search_query 전달)
                    # 기존 프롬프트의 {question} 자리에 위에서 만든 긴 지시사항이 들어갑니다.
                    response = st.session_state.chain.invoke({"query": search_query})
                    res = response['result']
                    
                except Exception as e:
                    res = f"Error: {str(e)}"
            
            # AI 응답 표시 및 저장
            chat_box.chat_message("assistant").write(res)
            st.session_state.messages.append({"role":"assistant", "content":res})

else:
    st.info(f"👈 {ui['sidebar_title']}")